#!/usr/bin/env python3
"""
rebuild_ppt.py -- ScratchV promotion PPT, fully rebuilt from scratch.

Key fixes over original:
- ZERO emoji unicode (all visual elements are native PPT shapes)
- ZERO U+000B vertical tab (all line breaks use \\n)
- Better layouts for slides 2, 9, 14, 15
- Clickable hyperlink on slide 18
- Rich shape-based visual design throughout
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ═══════════════════════════════════════════════════════════════
# Constants
# ═══════════════════════════════════════════════════════════════

FONT = "Microsoft YaHei"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.500)
LEFT_BAR_W = Inches(0.4)
TOP_BAR_H = Inches(0.06)
MARGIN = Inches(0.8)
CONTENT_X = MARGIN
CONTENT_W = Inches(11.733)

# ═══════════════════════════════════════════════════════════════
# Color System
# ═══════════════════════════════════════════════════════════════

C_DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)
C_DARK_CARD  = RGBColor(0x22, 0x22, 0x3E)
C_DARK_HOVER = RGBColor(0x2A, 0x2A, 0x42)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_WHITE_SOFT = RGBColor(0xF8, 0xFA, 0xFC)
C_TEXT_DARK  = RGBColor(0x2D, 0x2D, 0x2D)
C_TEXT_MUTED = RGBColor(0x66, 0x66, 0x66)
C_TEXT_LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
C_TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

C_TEAL       = RGBColor(0x00, 0x6D, 0x77)
C_TEAL_LIGHT = RGBColor(0xE0, 0xF7, 0xF7)
C_TEAL_PALE  = RGBColor(0xF0, 0xFA, 0xFA)
C_BLUE       = RGBColor(0x2E, 0x86, 0xAB)
C_BLUE_LIGHT = RGBColor(0xE8, 0xF4, 0xFA)
C_BLUE_PALE  = RGBColor(0xF4, 0xF8, 0xFC)
C_GOLD       = RGBColor(0xF5, 0xA6, 0x23)
C_GOLD_LIGHT = RGBColor(0xFF, 0xF8, 0xE8)
C_GREEN      = RGBColor(0x27, 0xAE, 0x60)
C_GREEN_LIGHT= RGBColor(0xEE, 0xFF, 0xF5)
C_ORANGE     = RGBColor(0xF3, 0x9C, 0x12)
C_ORANGE_LIGHT=RGBColor(0xFF, 0xF8, 0xEE)
C_RED        = RGBColor(0xE7, 0x4C, 0x3C)
C_RED_LIGHT  = RGBColor(0xFF, 0xF5, 0xF5)
C_RED_PALE   = RGBColor(0xFF, 0xF0, 0xEF)
C_BORDER     = RGBColor(0xE8, 0xE8, 0xE8)
C_BORDER_LT  = RGBColor(0xF0, 0xF0, 0xF0)

DIFF_COLORS  = {"低": C_GREEN, "中": C_ORANGE, "高": C_RED}
DIFF_LIGHTS  = {"低": C_GREEN_LIGHT, "中": C_ORANGE_LIGHT, "高": C_RED_LIGHT}
DIFF_PALES   = {"低": RGBColor(0xF2, 0xFC, 0xF5),
                "中": RGBColor(0xFF, 0xFB, 0xF2),
                "高": RGBColor(0xFF, 0xF5, 0xF4)}

# ═══════════════════════════════════════════════════════════════
# Helper Functions
# ═══════════════════════════════════════════════════════════════

def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def content_slide_decor(slide):
    """Standard content slide: left teal bar + top thin bar."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), LEFT_BAR_W, SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_TEAL
    bar.line.fill.background()

    top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, LEFT_BAR_W, Inches(0),
        SLIDE_W - LEFT_BAR_W, TOP_BAR_H)
    top.fill.solid()
    top.fill.fore_color.rgb = C_TEAL
    top.line.fill.background()


def add_decor_ovals(slide, positions, size=Inches(2.5), color=C_DARK_CARD):
    """Subtle decorative circles for dark background slides."""
    for x, y in positions:
        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), size, size)
        oval.fill.solid()
        oval.fill.fore_color.rgb = color
        oval.line.fill.background()


def add_slide_title(slide, title, subtitle=None, x=None, y=None, title_size=32):
    """Standard title block on content slides."""
    if x is None:
        x = CONTENT_X
    if y is None:
        y = Inches(0.5)
    txBox = slide.shapes.add_textbox(x, y, Inches(11.0), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.color.rgb = C_TEXT_DARK
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(15)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.font.bold = False
        p2.font.name = FONT
        p2.alignment = PP_ALIGN.LEFT
    return txBox


def add_ml_text(slide, left, top, width, height, text, size=14,
                color=C_TEXT_DARK, bold=False, align=PP_ALIGN.LEFT,
                line_spacing=1.3, italic=False, font_name=FONT):
    """Add multi-line text. \\n splits into separate paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line.strip()
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.font.italic = italic
        p.alignment = align
        if line_spacing:
            p.line_spacing = Pt(int(size * line_spacing))
    return txBox


def add_rounded_card(slide, left, top, width, height,
                     fill_color=C_WHITE, border_color=C_TEAL,
                     border_width=Pt(1.5)):
    """Standard rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_circle_badge(slide, left, top, diameter, fill_color, text,
                     font_size=11, text_color=C_WHITE):
    """Colored circle with centered number/label."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_difficulty_badge(slide, left, top, difficulty,
                         width=Inches(0.6), height=Inches(0.26)):
    """Colored difficulty badge."""
    c = DIFF_COLORS.get(difficulty, C_TEXT_MUTED)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = difficulty
    p.font.size = Pt(9)
    p.font.color.rgb = C_WHITE
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_arrow_right(slide, left, top, width, height, color):
    """Right-pointing arrow for flow diagrams."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_arrow_down(slide, left, top, width, height, color):
    """Down-pointing arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_accent_bar(slide, left, top, width, height, color):
    """Thin accent bar for card headers or visual flair."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_chevron(slide, left, top, width, height, color):
    """Chevron shape for process flows."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_section_divider(slide, left, top, width, color, text=None):
    """Horizontal section divider line with optional label."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.02))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    if text:
        add_ml_text(slide, left, top + Inches(0.08), width, Inches(0.3),
                    text, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_small_icon_circle(slide, left, top, diameter, color, char, char_size=12):
    """Small colored circle with a single character (like bullet or check)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = char
    p.font.size = Pt(char_size)
    p.font.color.rgb = C_WHITE
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Cover / Hero
# ═══════════════════════════════════════════════════════════════

def build_slide_01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK_BG)

    # Decorative geometric elements -- layered ovals on right
    # Largest (bottom layer, naturally clipped by slide edge)
    add_decor_ovals(slide, [(10.0, 4.5)], Inches(3.0), C_DARK_CARD)
    # Medium
    add_decor_ovals(slide, [(11.2, 5.3)], Inches(2.0), C_DARK_HOVER)
    # Small accent dots
    for (dx, dy) in [(11.8, 2.0), (12.5, 3.8), (12.8, 1.0)]:
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(dx), Inches(dy), Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_TEAL
        dot.line.fill.background()

    # Left side subtle decoration
    add_decor_ovals(slide, [(-0.8, 4.5)], Inches(2.0), C_DARK_CARD)

    # Top accent line
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), C_TEAL)

    # Main title
    add_ml_text(slide, Inches(1.2), Inches(1.0), Inches(10.0), Inches(1.4),
                "探索「AI 模型 -> 芯片指令」的\n神奇之旅",
                size=46, color=C_TEXT_WHITE, bold=True, line_spacing=1.4)

    # Subtitle
    add_ml_text(slide, Inches(1.2), Inches(3.0), Inches(10.0), Inches(0.5),
                "零基础友好的 AI 编译器开源项目 | 线上宣讲",
                size=24, color=C_TEXT_LIGHT)

    # Info row 1 -- date
    add_circle_badge(slide, Inches(1.2), Inches(4.2), Inches(0.35),
                     C_TEAL, "1", 12, C_WHITE)
    add_ml_text(slide, Inches(1.7), Inches(4.15), Inches(8.0), Inches(0.4),
                "2025 年 6 月  |  三个月，从零搭建你的第一个编译器",
                size=16, color=C_TEXT_LIGHT)

    # Info row 2 -- no prerequisites
    add_circle_badge(slide, Inches(1.2), Inches(4.8), Inches(0.35),
                     C_GOLD, "i", 13, C_WHITE)
    add_ml_text(slide, Inches(1.7), Inches(4.75), Inches(10.0), Inches(0.4),
                "不需要编译原理基础 | 不需要 AI 基础 | 只需要好奇心和耐心",
                size=14, color=C_TEXT_MUTED)

    # Bottom gold line
    add_accent_bar(slide, Inches(1.2), Inches(5.8), Inches(4.0), Inches(0.03), C_GOLD)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: Agenda -- 7 modules in 2-column card grid
# ═══════════════════════════════════════════════════════════════

def build_slide_02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "宣讲提纲",
                    "七个模块，带你从零走进编译器的世界")

    modules = [
        ("01", "项目概览", "我们要做什么？", C_TEAL),
        ("02", "为什么值得参加", "你能收获什么", C_BLUE),
        ("03", "三个月学习路线", "Phase 1-2-3 逐级进阶", C_GREEN),
        ("04", "课题精选", "14 个课题全方位解析", C_GOLD),
        ("05", "时间节点与适合人群", "关键里程碑 & 报名要求", C_ORANGE),
        ("06", "如何上手", "参与准备与上手路线", C_TEAL),
        ("07", "Q & A", "常见疑问与报名方式", C_BLUE),
    ]

    # 2-column grid: left 4, right 3
    card_w = Inches(5.35)
    card_h = Inches(1.1)
    gap_x = Inches(0.5)
    gap_y = Inches(0.18)
    total_w = 2 * card_w + gap_x
    start_x = (SLIDE_W - total_w) / 2
    start_y = Inches(1.65)

    left_col = modules[:4]
    right_col = modules[4:]

    for col_idx, col_modules in enumerate([left_col, right_col]):
        for row_idx, (num, title, subtitle, accent) in enumerate(col_modules):
            x = start_x + col_idx * (card_w + gap_x)
            y = start_y + row_idx * (card_h + gap_y)

            # Card background
            add_rounded_card(slide, x, y, card_w, card_h,
                           C_WHITE, accent, Pt(1.2))
            # Left accent strip
            add_accent_bar(slide, x, y, Inches(0.07), card_h, accent)

            # Number circle
            add_circle_badge(slide, x + Inches(0.3), y + Inches(0.22),
                           Inches(0.5), accent, num, 14, C_WHITE)

            # Title
            add_ml_text(slide, x + Inches(1.0), y + Inches(0.15),
                       card_w - Inches(1.3), Inches(0.35),
                       title, size=17, color=C_TEXT_DARK, bold=True)
            # Subtitle
            add_ml_text(slide, x + Inches(1.0), y + Inches(0.55),
                       card_w - Inches(1.3), Inches(0.3),
                       subtitle, size=12, color=C_TEXT_MUTED)

    # Bottom connecting note
    add_section_divider(slide, start_x, Inches(6.8), total_w, C_TEAL,
                        "全程零基础友好 -- 每步都有 mentor 指导")


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: Curiosity Hook -- 3 question cards
# ═══════════════════════════════════════════════════════════════

def build_slide_03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "你有没有好奇过......",
                    "一些可能让你「心里痒痒」的问题")

    card_w = Inches(3.8)
    card_h = Inches(4.0)
    gap = (SLIDE_W - 3 * card_w - MARGIN * 2) / 2
    cols = [MARGIN, MARGIN + card_w + gap, MARGIN + 2 * (card_w + gap)]
    card_y = Inches(1.8)

    cards_data = [
        (C_TEAL, C_TEAL_LIGHT,
         "你写的 Python 代码，\n电脑到底是怎么\n「听懂」并执行的？",
         "从高级语言到机器指令，\n中间经历了怎样的魔法转换？"),
        (C_BLUE, C_BLUE_LIGHT,
         "那些炫酷的 AI 模型，\n最终怎么在小小的\n芯片上跑起来？",
         "模型是数学公式，芯片只懂 0 和 1 --\n谁来当这个翻译官？"),
        (C_GOLD, C_GOLD_LIGHT,
         "编译器 --\n这个听起来很高深的东西，\n到底在做什么？",
         "它不是一个黑盒，\n而是你可以亲手搭建的工具。"),
    ]

    for i, (accent, fill_c, q_text, d_text) in enumerate(cards_data):
        x = cols[i]
        # Card
        add_rounded_card(slide, x, card_y, card_w, card_h,
                        fill_color=fill_c, border_color=accent, border_width=Pt(1.5))
        # Top accent bar
        add_accent_bar(slide, x, card_y, card_w, Inches(0.06), accent)

        # Decorative question mark circle (replacing emoji)
        add_circle_badge(slide, x + card_w / 2 - Inches(0.35),
                        card_y + Inches(0.3), Inches(0.7),
                        accent, "?", 24, C_WHITE)

        # Question text
        add_ml_text(slide, x + Inches(0.3), card_y + Inches(1.25),
                    card_w - Inches(0.6), Inches(1.5),
                    q_text, size=14, color=C_TEXT_DARK, bold=True,
                    line_spacing=1.4)
        # Separator
        add_accent_bar(slide, x + Inches(0.4), card_y + Inches(2.85),
                       card_w - Inches(0.8), Inches(0.015), accent)
        # Description
        add_ml_text(slide, x + Inches(0.3), card_y + Inches(3.0),
                    card_w - Inches(0.6), Inches(0.9),
                    d_text, size=12, color=C_TEXT_MUTED, line_spacing=1.35)

    # Bottom CTA
    add_ml_text(slide, Inches(2.0), Inches(6.1), Inches(9.0), Inches(0.5),
                "-- 哪怕你只学过一点点编程，这个项目就是为你准备的 --",
                size=18, color=C_TEAL, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: What We Build -- Input -> Process -> Output
# ═══════════════════════════════════════════════════════════════

def build_slide_04(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "我们要做什么？",
                    "用三个月时间，从零搭建一个迷你 AI 编译器")

    card_w = Inches(3.5)
    card_h = Inches(3.0)
    total_w = 3 * card_w + 2 * Inches(0.4)
    start_x = (SLIDE_W - total_w) / 2
    card_y = Inches(1.9)
    arrow_w = Inches(0.4)

    stages = [
        ("Input", "一个简单的 ONNX\nAI 模型文件\n（加法和乘法）", C_TEAL, "in"),
        ("Compiler Core", "[1] 读懂模型\n[2] 翻译成中间语言\n[3] 优化\n[4] 生成指令", C_BLUE, "core"),
        ("Output", "RISC-V 汇编代码\n（add, load 等\n芯片指令）", C_GOLD, "out"),
    ]

    for i, (label, desc, accent, tag) in enumerate(stages):
        x = start_x + i * (card_w + arrow_w)
        # Card
        add_rounded_card(slide, x, card_y, card_w, card_h, C_WHITE, accent, Pt(1.5))
        # Top accent bar
        add_accent_bar(slide, x, card_y, card_w, Inches(0.06), accent)
        # Label
        add_ml_text(slide, x, card_y + Inches(0.2), card_w, Inches(0.4),
                    label, size=18, color=accent, bold=True, align=PP_ALIGN.CENTER)
        # Label underline
        add_accent_bar(slide, x + Inches(0.5), card_y + Inches(0.65),
                       card_w - Inches(1.0), Inches(0.02), accent)
        # Description
        add_ml_text(slide, x + Inches(0.3), card_y + Inches(0.85),
                    card_w - Inches(0.6), Inches(1.8),
                    desc, size=13, color=C_TEXT_DARK, align=PP_ALIGN.CENTER,
                    line_spacing=1.45)

        # Arrow between cards
        if i < 2:
            ax = x + card_w + Inches(0.02)
            add_arrow_right(slide, ax, card_y + Inches(1.35),
                           arrow_w - Inches(0.04), Inches(0.25), C_GOLD)

    # Bottom feature checkmarks
    features = [
        "完全自主实现，不依赖 LLVM / MLIR 等巨型框架",
        "放到模拟器 (tinyfive) 中运行验证",
        "每一步都亲手写出来，真正搞懂原理",
    ]
    for i, feat in enumerate(features):
        y = Inches(5.3) + i * Inches(0.40)
        add_small_icon_circle(slide, Inches(2.3), y + Inches(0.02),
                             Inches(0.28), C_GREEN, "v", 10)
        add_ml_text(slide, Inches(2.75), y, Inches(8.0), Inches(0.35),
                    feat, size=13, color=C_TEXT_DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Full Pipeline -- 6 stages
# ═══════════════════════════════════════════════════════════════

def build_slide_05(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "从模型到芯片的完整链路",
                    "编译器 -- AI 世界与芯片世界的桥梁")

    stages = [
        ("ONNX 模型", "读取并解析\nAI 模型文件", C_DARK_BG),
        ("IR 中间表示", "翻译成自己\n定义的指令", C_TEAL),
        ("优化器", "死代码消除\n常量折叠等", C_DARK_BG),
        ("后端", "指令选择\n寄存器分配", C_TEAL),
        ("RISC-V 汇编", "生成芯片\n可执行指令", C_DARK_BG),
        ("模拟器 tinyfive", "运行验证\n输出结果", C_TEAL),
    ]

    n = len(stages)
    stage_w = Inches(1.9)
    arrow_w = Inches(0.18)
    total_w = n * stage_w + (n - 1) * arrow_w
    start_x = (SLIDE_W - total_w) / 2
    card_y = Inches(2.1)
    card_h_title = Inches(0.5)
    card_h_desc = Inches(1.3)

    for i, (name, desc, accent) in enumerate(stages):
        x = start_x + i * (stage_w + arrow_w)
        # Title bar
        title_bar = add_rounded_card(slide, x, card_y, stage_w, card_h_title,
                                     fill_color=accent, border_color=None)
        text_color = C_WHITE if accent in (C_DARK_BG, C_TEAL) else C_TEXT_WHITE
        add_ml_text(slide, x, card_y + Inches(0.02), stage_w, card_h_title,
                    name, size=10, color=C_TEXT_WHITE, bold=True,
                    align=PP_ALIGN.CENTER, line_spacing=1.0)
        # Description
        desc_card = add_rounded_card(slide, x, card_y + card_h_title,
                                     stage_w, card_h_desc,
                                     fill_color=C_WHITE_SOFT, border_color=accent,
                                     border_width=Pt(1))
        add_ml_text(slide, x + Inches(0.08), card_y + card_h_title + Inches(0.12),
                    stage_w - Inches(0.16), card_h_desc - Inches(0.2),
                    desc, size=9, color=C_TEXT_DARK, align=PP_ALIGN.CENTER,
                    line_spacing=1.4)

        # Arrow between stages
        if i < n - 1:
            ax = x + stage_w
            add_arrow_right(slide, ax + Inches(0.01), card_y + Inches(0.1),
                           arrow_w - Inches(0.02), Inches(0.2), C_GOLD)

    # Bottom insight
    y_note = card_y + card_h_title + card_h_desc + Inches(0.45)
    add_small_icon_circle(slide, Inches(1.5), y_note, Inches(0.3), C_TEAL, "i", 11)
    add_ml_text(slide, Inches(1.95), y_note - Inches(0.02), Inches(9.5), Inches(0.4),
                "整个流程没有任何黑盒 -- 从模型文件的第一行，到汇编指令的最后一行，都由你亲手构建",
                size=14, color=C_TEXT_DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Why Participate -- 3 reason cards
# ═══════════════════════════════════════════════════════════════

def build_slide_06(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "为什么值得你来试一试？",
                    "三个理由，总有一个打动你")

    card_w = Inches(3.8)
    card_h = Inches(5.2)
    gap = (SLIDE_W - 3 * card_w - MARGIN * 2) / 2
    cols = [MARGIN, MARGIN + card_w + gap, MARGIN + 2 * (card_w + gap)]
    card_y = Inches(1.7)

    reasons = [
        (C_TEAL, "1", "不需要大神基础",
         "学过一点 Python 或 C？够了。\n知道数组、函数、循环？完全够。\n"
         "两周带你入门 RISC-V，\n不需要预先掌握任何底层知识。",
         "像搭积木一样，一块一块搭起来"),
        (C_BLUE, "2", "真正理解底层",
         "不再是只会 import torch 的调包侠。\n"
         "理解从数学模型到机器指令的完整链路。\n"
         "高性能计算、AI 芯片的核心能力，\n从这里开始培养。",
         "简历加分：独立实现 AI -> RISC-V 完整编译器"),
        (C_GOLD, "3", "温暖的开源社区",
         "每周线上答疑 + 讲解。\n"
         "Peer Review 代码，互相改 bug。\n"
         "一起庆祝每个里程碑的达成。\n"
         "你不是一个人在战斗。",
         "你的代码会成为开源项目的一部分"),
    ]

    for i, (accent, num, title, desc, badge_text) in enumerate(reasons):
        x = cols[i]
        # Card
        add_rounded_card(slide, x, card_y, card_w, card_h,
                        C_WHITE, accent, Pt(1.5))
        # Top accent bar
        add_accent_bar(slide, x, card_y, card_w, Inches(0.06), accent)

        # Number circle at top-center
        add_circle_badge(slide, x + card_w / 2 - Inches(0.3),
                        card_y + Inches(0.3), Inches(0.6),
                        accent, num, 18, C_WHITE)

        # Title
        add_ml_text(slide, x + Inches(0.3), card_y + Inches(1.1),
                    card_w - Inches(0.6), Inches(0.4),
                    title, size=19, color=accent, bold=True,
                    align=PP_ALIGN.CENTER)

        # Title underline
        add_accent_bar(slide, x + Inches(0.8), card_y + Inches(1.5),
                       card_w - Inches(1.6), Inches(0.02), accent)

        # Description
        add_ml_text(slide, x + Inches(0.35), card_y + Inches(1.7),
                    card_w - Inches(0.7), Inches(2.4),
                    desc, size=12, color=C_TEXT_DARK, line_spacing=1.55)

        # Bottom badge -- accent bar style, not filling entire bottom
        badge_y = card_y + card_h - Inches(0.6)
        add_accent_bar(slide, x, badge_y, card_w, Inches(0.03), accent)
        add_ml_text(slide, x + Inches(0.2), badge_y + Inches(0.08),
                    card_w - Inches(0.4), Inches(0.45),
                    badge_text, size=11, color=accent, bold=True,
                    align=PP_ALIGN.CENTER, italic=True, line_spacing=1.2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: Roadmap -- 3 Phase cards
# ═══════════════════════════════════════════════════════════════

def build_slide_07(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "三个月学习路线",
                    "循序渐进的三阶段学习计划，每个阶段都有明确目标")

    card_w = Inches(3.8)
    card_h = Inches(5.0)
    gap = (SLIDE_W - 3 * card_w - MARGIN * 2) / 2
    cols = [MARGIN, MARGIN + card_w + gap, MARGIN + 2 * (card_w + gap)]
    card_y = Inches(1.7)

    phases = [
        (C_TEAL, "Phase 1", "第 1 - 2 周", "跑通完整链路",
         "跑通别人写好的完整例子\n看懂「模型 -> 指令」的全过程\n建立对编译器的整体认知",
         "哇，原来是这样！"),
        (C_BLUE, "Phase 2", "第 3 - 8 周", "亲手实现编译器",
         "把 ONNX 模型翻译成自己的 IR\n实现 RISC-V 后端生成汇编\n实现基础优化（常量折叠等）",
         "开始创造，成就感爆棚"),
        (C_GOLD, "Phase 3", "第 9 - 12 周", "优化与产出",
         "跑通更多模型\n指令更短、运行更快\n编写文档，完善项目",
         "我居然做出了一个编译器！"),
    ]

    for i, (accent, phase, weeks, title, items, reaction) in enumerate(phases):
        x = cols[i]
        # Card
        add_rounded_card(slide, x, card_y, card_w, card_h, C_WHITE, accent, Pt(1.5))

        # Header bar (taller for clarity)
        header_h = Inches(0.9)
        header = add_rounded_card(slide, x, card_y, card_w, header_h,
                                  fill_color=accent, border_color=None)
        add_ml_text(slide, x + Inches(0.2), card_y + Inches(0.05),
                    card_w - Inches(0.4), Inches(0.35),
                    phase, size=16, color=C_WHITE, bold=True,
                    align=PP_ALIGN.CENTER)
        add_ml_text(slide, x + Inches(0.2), card_y + Inches(0.45),
                    card_w - Inches(0.4), Inches(0.25),
                    weeks, size=11, color=C_WHITE,
                    align=PP_ALIGN.CENTER)

        # Section title inside card
        add_ml_text(slide, x + Inches(0.3), card_y + header_h + Inches(0.2),
                    card_w - Inches(0.6), Inches(0.35),
                    title, size=16, color=C_TEXT_DARK, bold=True,
                    align=PP_ALIGN.CENTER)

        # Accent line under title
        add_accent_bar(slide, x + Inches(0.6), card_y + header_h + Inches(0.55),
                       card_w - Inches(1.2), Inches(0.02), accent)

        # Checklist items
        for j, item in enumerate(items.split('\n')):
            iy = card_y + header_h + Inches(0.75) + j * Inches(0.45)
            add_small_icon_circle(slide, x + Inches(0.3), iy + Inches(0.02),
                                 Inches(0.22), accent, "v", 8)
            add_ml_text(slide, x + Inches(0.65), iy,
                       card_w - Inches(0.9), Inches(0.35),
                       item, size=12, color=C_TEXT_DARK, line_spacing=1.0)

        # Reaction at bottom
        react_y = card_y + card_h - Inches(0.65)
        add_accent_bar(slide, x + Inches(0.4), react_y,
                       card_w - Inches(0.8), Inches(0.02), accent)
        add_ml_text(slide, x + Inches(0.2), react_y + Inches(0.08),
                    card_w - Inches(0.4), Inches(0.45),
                    reaction, size=14, color=accent, bold=False,
                    align=PP_ALIGN.CENTER, italic=True)

    # Bottom note
    note_y = Inches(6.85)
    add_section_divider(slide, MARGIN, note_y, CONTENT_W, C_TEAL,
                        "每周只需 8 - 10 小时  |  提供预置框架和 benchmark  |  每周线上答疑")


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: Topics Overview -- 3-column by difficulty
# ═══════════════════════════════════════════════════════════════

def build_slide_08(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "课题精选总览",
                    "14 个课题覆盖编译器的各个模块，按难度分层")

    topics_by_diff = {
        "低": [
            ("7", "编译器日志增强器"),
            ("13", "窥孔优化器"),
            ("14", "常量加载合并优化"),
            ("5", "RISC-V 汇编代码美化器"),
            ("20", "项目代码规范与格式化"),
        ],
        "中": [
            ("1", "DSL 前端增强器"),
            ("6", "编译器性能测试套件"),
            ("9", "DSL 错误提示美化器"),
            ("21", "IR 验证器"),
            ("28", "完善后端指令选择"),
        ],
        "高": [
            ("11", "控制流图（CFG）生成器"),
            ("12", "RISC-V 后端指令计数统计器"),
            ("17", "寄存器分配（线性扫描）"),
            ("18", "指令调度（列表调度）"),
        ],
    }

    col_labels = [
        ("低", "入门友好", C_GREEN),
        ("中", "进阶挑战", C_ORANGE),
        ("高", "大神修炼", C_RED),
    ]
    col_w = Inches(3.7)
    gap = (SLIDE_W - 3 * col_w - MARGIN * 2) / 2
    col_x = [MARGIN, MARGIN + col_w + gap, MARGIN + 2 * (col_w + gap)]
    start_y = Inches(1.8)

    for ci, (diff, label, color) in enumerate(col_labels):
        x = col_x[ci]
        # Column header
        header = add_rounded_card(slide, x, start_y, col_w, Inches(0.42),
                                  fill_color=color, border_color=None)
        add_ml_text(slide, x + Inches(0.1), start_y + Inches(0.02),
                    col_w - Inches(0.2), Inches(0.38),
                    f"[{diff}]  {label}", size=14, color=C_WHITE, bold=True,
                    align=PP_ALIGN.CENTER)

        # Topic items
        topics = topics_by_diff[diff]
        for ti, (num, name) in enumerate(topics):
            y = start_y + Inches(0.58) + ti * Inches(0.65)
            # Number circle
            add_circle_badge(slide, x + Inches(0.2), y + Inches(0.08),
                           Inches(0.34), color, num, 10, C_WHITE)
            # Topic name
            add_ml_text(slide, x + Inches(0.7), y + Inches(0.02),
                        col_w - Inches(0.9), Inches(0.35),
                        name, size=13, color=C_TEXT_DARK, bold=False)
            # Separator
            if ti < len(topics) - 1:
                add_accent_bar(slide, x + Inches(0.3), y + Inches(0.52),
                               col_w - Inches(0.6), Inches(0.008), C_BORDER_LT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: Beginner Topics -- 5 cards in 3+2 centered layout
# ═══════════════════════════════════════════════════════════════

def build_slide_09(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "课题详解 | 入门友好型",
                    "适合新手起步，快速获得成就感  [ 难度：低 ]")

    topics = [
        ("#7", "编译器日志增强器",
         "为编译器添加结构化日志输出\n支持不同日志级别 (INFO/WARN/ERROR)\n帮助开发者调试编译流程"),
        ("#13", "窥孔优化器",
         "扫描生成的汇编代码\n用更短的指令序列替换冗余模式\n经典的局部优化技术"),
        ("#14", "常量加载合并优化",
         "识别重复的常量加载指令\n合并为一次加载，减少代码体积\n简单但效果明显的优化"),
        ("#5", "RISC-V 汇编代码美化器",
         "格式化汇编输出\n添加注释和可读性改进\n让生成的代码更易读"),
        ("#20", "项目代码规范与格式化",
         "统一代码风格\n配置 linter / formatter\n提升项目的专业性"),
    ]

    card_w = Inches(3.6)
    card_h = Inches(2.1)
    gap_x = Inches(0.3)
    gap_y = Inches(0.28)
    # Row 1: 3 cards, Row 2: 2 cards -- both CENTERED
    row1_w = 3 * card_w + 2 * gap_x
    row2_w = 2 * card_w + 1 * gap_x
    row1_start_x = (SLIDE_W - row1_w) / 2
    row2_start_x = (SLIDE_W - row2_w) / 2
    start_y = Inches(1.85)

    for i, (num, title, desc) in enumerate(topics):
        if i < 3:
            x = row1_start_x + i * (card_w + gap_x)
            y = start_y
        else:
            x = row2_start_x + (i - 3) * (card_w + gap_x)
            y = start_y + card_h + gap_y

        # Card
        add_rounded_card(slide, x, y, card_w, card_h, C_WHITE, C_GREEN, Pt(1.2))
        # Left accent strip (wider for visibility)
        add_accent_bar(slide, x, y, Inches(0.08), card_h, C_GREEN)

        # Number
        add_ml_text(slide, x + Inches(0.25), y + Inches(0.15),
                    Inches(1.5), Inches(0.3),
                    num, size=14, color=C_GREEN, bold=True)
        # Difficulty badge
        add_difficulty_badge(slide, x + card_w - Inches(0.85), y + Inches(0.17), "低")

        # Title
        add_ml_text(slide, x + Inches(0.25), y + Inches(0.55),
                    card_w - Inches(0.5), Inches(0.3),
                    title, size=14, color=C_TEXT_DARK, bold=True)
        # Description
        add_ml_text(slide, x + Inches(0.25), y + Inches(0.95),
                    card_w - Inches(0.5), Inches(1.0),
                    desc, size=10, color=C_TEXT_MUTED, line_spacing=1.4)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: Intermediate Topics -- 5 cards in grid
# ═══════════════════════════════════════════════════════════════

def build_slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "课题详解 | 进阶挑战型",
                    "需要一定的编程能力，深入编译器核心  [ 难度：中 ]")

    topics = [
        ("#1", "DSL 前端增强器",
         "扩展 DSL 语法支持\n支持更多算子类型\n改进前端解析能力"),
        ("#9", "DSL 错误提示美化器",
         "友好的错误信息展示\n精确的错误位置定位\n彩色终端输出"),
        ("#6", "编译器性能测试套件",
         "自动化 benchmark 流程\n统计编译时间和代码质量\n生成性能对比报告"),
        ("#21", "IR 验证器",
         "检查 IR 的合法性\n验证数据流正确性\n确保优化不改变语义"),
        ("#28", "完善后端指令选择",
         "扩展指令选择规则\n支持更多 RISC-V 指令\n后端功能完善"),
    ]

    card_w = Inches(3.6)
    card_h = Inches(2.1)
    gap_x = Inches(0.3)
    gap_y = Inches(0.28)
    cols_per_row = 3

    for idx, (num, title, desc) in enumerate(topics):
        row = idx // cols_per_row
        col = idx % cols_per_row
        items_in_row = min(cols_per_row, len(topics) - row * cols_per_row)
        row_w = items_in_row * card_w + (items_in_row - 1) * gap_x
        start_x = (SLIDE_W - row_w) / 2
        x = start_x + col * (card_w + gap_x)
        y = Inches(1.85) + row * (card_h + gap_y)

        add_rounded_card(slide, x, y, card_w, card_h, C_WHITE, C_ORANGE, Pt(1.2))
        add_accent_bar(slide, x, y, Inches(0.08), card_h, C_ORANGE)
        add_ml_text(slide, x + Inches(0.25), y + Inches(0.15),
                    Inches(1.5), Inches(0.3),
                    num, size=14, color=C_ORANGE, bold=True)
        add_difficulty_badge(slide, x + card_w - Inches(0.85), y + Inches(0.17), "中")
        add_ml_text(slide, x + Inches(0.25), y + Inches(0.55),
                    card_w - Inches(0.5), Inches(0.3),
                    title, size=14, color=C_TEXT_DARK, bold=True)
        add_ml_text(slide, x + Inches(0.25), y + Inches(0.95),
                    card_w - Inches(0.5), Inches(1.0),
                    desc, size=10, color=C_TEXT_MUTED, line_spacing=1.4)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: Advanced Topics -- 4 cards in a row
# ═══════════════════════════════════════════════════════════════

def build_slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "课题详解 | 大神修炼型",
                    "挑战编译器核心算法，收获满满硬技能  [ 难度：高 ]")

    topics = [
        ("#11", "控制流图(CFG)生成器",
         "从 IR 中构建基本块\n生成控制流图可视化\n为后续优化奠定基础",
         "编译优化 | 图算法"),
        ("#12", "后端指令计数统计器",
         "统计生成的各类指令数量\n分析指令分布特征\n输出可视化统计报告",
         "数据分析 | 可视化"),
        ("#17", "寄存器分配(线性扫描)",
         "实现寄存器分配算法\n处理变量生命周期\n优化寄存器使用效率",
         "经典算法实现"),
        ("#18", "指令调度(列表调度)",
         "分析指令依赖关系\n重排指令顺序\n提高流水线效率",
         "编译优化 | 调度算法"),
    ]

    card_w = Inches(2.95)
    card_h = Inches(4.0)
    total_w = 4 * card_w + 3 * Inches(0.2)
    gap = (SLIDE_W - total_w - 2 * MARGIN) / 2 if total_w + 2 * MARGIN < SLIDE_W else Inches(0.2)
    start_x = (SLIDE_W - total_w) / 2
    start_y = Inches(1.9)

    for i, (num, title, desc, tag) in enumerate(topics):
        x = start_x + i * (card_w + Inches(0.2))
        y = start_y

        add_rounded_card(slide, x, y, card_w, card_h, C_RED_PALE, C_RED, Pt(1.3))
        # Top accent
        add_accent_bar(slide, x, y, card_w, Inches(0.06), C_RED)
        # Number + difficulty
        add_ml_text(slide, x + Inches(0.2), y + Inches(0.2),
                    Inches(1.0), Inches(0.3),
                    num, size=15, color=C_RED, bold=True)
        add_difficulty_badge(slide, x + card_w - Inches(0.85), y + Inches(0.22), "高")
        # Title
        add_ml_text(slide, x + Inches(0.2), y + Inches(0.65),
                    card_w - Inches(0.4), Inches(0.5),
                    title, size=13, color=C_TEXT_DARK, bold=True,
                    line_spacing=1.25)
        # Accent line
        add_accent_bar(slide, x + Inches(0.2), y + Inches(1.15),
                       card_w - Inches(0.4), Inches(0.015), C_RED)
        # Description
        add_ml_text(slide, x + Inches(0.2), y + Inches(1.3),
                    card_w - Inches(0.4), Inches(1.5),
                    desc, size=11, color=C_TEXT_MUTED, line_spacing=1.5)
        # Bottom tag
        tag_bar = add_rounded_card(slide, x + Inches(0.2), y + card_h - Inches(0.55),
                                   card_w - Inches(0.4), Inches(0.35),
                                   fill_color=C_RED, border_color=None)
        add_ml_text(slide, x + Inches(0.2), y + card_h - Inches(0.5),
                    card_w - Inches(0.4), Inches(0.3),
                    tag, size=9, color=C_TEXT_WHITE, bold=True,
                    align=PP_ALIGN.CENTER)

    # Bottom reassurance
    add_ml_text(slide, Inches(1.0), Inches(6.2), Inches(11.0), Inches(0.5),
                "高难度课题有 mentor 全程指导，不用担心做不出来",
                size=14, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12: Eligibility -- checkmark bars
# ═══════════════════════════════════════════════════════════════

def build_slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "谁适合报名？")

    items = [
        "大二、大三、研一同学，或自学编程爱好者",
        "学过一门编程语言（Python / C / C++ 都行）",
        "对「计算机到底怎么跑程序」有好奇心",
        "不怕犯错，敢写代码（Bug 是学习的一部分！）",
        "每周能拿出 8 - 10 小时",
    ]

    bar_w = Inches(10.5)
    bar_h = Inches(0.7)
    start_x = Inches(1.5)
    start_y = Inches(1.9)
    gap = Inches(0.15)

    for i, item in enumerate(items):
        y = start_y + i * (bar_h + gap)
        add_rounded_card(slide, start_x, y, bar_w, bar_h,
                        C_GREEN_LIGHT, C_GREEN, Pt(1))
        add_small_icon_circle(slide, start_x + Inches(0.2), y + Inches(0.16),
                             Inches(0.38), C_GREEN, "v", 12)
        add_ml_text(slide, start_x + Inches(0.75), y + Inches(0.1),
                    bar_w - Inches(1.0), bar_h - Inches(0.15),
                    item, size=16, color=C_TEXT_DARK)

    # Encouragement
    y_enc = start_y + len(items) * (bar_h + gap) + Inches(0.5)
    add_ml_text(slide, Inches(1.5), y_enc, Inches(10.0), Inches(0.7),
                "你可能还没学过编译原理、还没搞懂指令集、甚至对汇编有点畏惧 -- 都没关系。\n"
                "我们就是来带你一步步跨过这些坎的。",
                size=15, color=C_TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13: Timeline -- horizontal milestone markers
# ═══════════════════════════════════════════════════════════════

def build_slide_13(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "关键时间节点",
                    "从报名到结项，每一步都有清晰的时间线")

    milestones = [
        ("即日起", "开始报名", C_TEAL),
        ("6 / 20", "线上宣讲\n课题选择\n报名截止", C_TEAL),
        ("7 / 10", "项目正式开启\n启动会\n第一周任务", C_TEAL),
        ("8 / 1", "Phase 1 验收\n第一个里程碑", C_TEAL),
        ("8 / 28", "Phase 2 验收\n第二个里程碑", C_TEAL),
        ("9 / 27", "项目结项\n成果展示+总结", C_GOLD),
    ]

    n = len(milestones)
    line_y = Inches(3.2)
    node_spacing = Inches(2.05)
    total_line_w = (n - 1) * node_spacing
    line_start_x = (SLIDE_W - total_line_w) / 2
    node_radius = Inches(0.2)

    # Horizontal timeline line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, line_start_x, line_y - Inches(0.02),
        total_line_w, Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = C_TEAL
    line.line.fill.background()

    for i, (date, event, color) in enumerate(milestones):
        x = line_start_x + i * node_spacing

        # Node circle
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x - node_radius, line_y - node_radius,
            node_radius * 2, node_radius * 2)
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()
        # White inner dot
        inner = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x - Inches(0.06), line_y - Inches(0.06),
            Inches(0.12), Inches(0.12))
        inner.fill.solid()
        inner.fill.fore_color.rgb = C_WHITE
        inner.line.fill.background()

        # Date above
        add_ml_text(slide, x - Inches(0.7), line_y - Inches(0.75),
                    Inches(1.4), Inches(0.35),
                    date, size=14, color=color, bold=True,
                    align=PP_ALIGN.CENTER)
        # Event below
        add_ml_text(slide, x - Inches(0.85), line_y + Inches(0.35),
                    Inches(1.7), Inches(1.2),
                    event, size=10, color=C_TEXT_DARK,
                    align=PP_ALIGN.CENTER, line_spacing=1.35)

    # Bottom note
    add_ml_text(slide, Inches(1.5), Inches(5.8), Inches(10.0), Inches(0.5),
                "每个里程碑都有进度检查和 mentor 指导，确保你不会掉队",
                size=14, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14: Development Workflow -- left flowchart + right support
# ═══════════════════════════════════════════════════════════════

def build_slide_14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "开发流程",
                    "从选择课题到项目结项，完整的参与流程与支持体系")

    # ── LEFT: 7-step vertical flowchart ──
    step_x = Inches(0.9)
    step_w = Inches(5.9)
    step_h = Inches(0.58)
    step_start_y = Inches(1.7)
    step_gap = Inches(0.06)
    circle_d = Inches(0.4)

    steps = [
        "选择课题、立项",
        "完成设计文档",
        "编写程序与对应单元测试",
        "提交代码，运行并查看结果",
        "效果符合预期后提交 PR 并合入",
        "完善开发文档、实验报告",
        "结项",
    ]

    for i, step_text in enumerate(steps):
        y = step_start_y + i * (step_h + step_gap)
        # Step card
        add_rounded_card(slide, step_x, y, step_w, step_h,
                        C_WHITE, C_TEAL, Pt(1.2))
        # Number circle
        add_circle_badge(slide, step_x + Inches(0.12), y + Inches(0.09),
                        circle_d, C_TEAL, str(i + 1), 11, C_WHITE)
        # Step text
        add_ml_text(slide, step_x + Inches(0.65), y + Inches(0.08),
                    step_w - Inches(0.9), step_h - Inches(0.1),
                    step_text, size=13, color=C_TEXT_DARK)

        # Down arrow between steps
        if i < len(steps) - 1:
            ay = y + step_h
            add_arrow_down(slide, step_x + step_w / 2 - Inches(0.08), ay,
                          Inches(0.16), step_gap - Inches(0.01), C_TEAL)

    # ── RIGHT: Support cards ──
    support_x = Inches(7.3)
    support_w = Inches(5.3)
    support_start_y = Inches(1.7)

    support_cards = [
        ("定期答疑", "每周固定时间段内答疑",
         "遇到问题随时获得帮助，\n不会卡住超过一周。\nmentor 主动跟进你的进度。"),
        ("进度同步", "文档形式同步进度",
         "每周汇报进展，\nmentor 及时掌握你的状态，\n确保方向正确。"),
        ("小组协作", "多个课题组成小组",
         "方便随时交流问题，\n互相 review 代码，\n互相打气，共同成长。"),
    ]

    scard_h = Inches(1.55)
    sgap = Inches(0.18)

    for i, (title, subtitle, desc) in enumerate(support_cards):
        y = support_start_y + i * (scard_h + sgap)
        add_rounded_card(slide, support_x, y, support_w, scard_h,
                        C_BLUE_LIGHT, C_BLUE, Pt(1.2))
        add_accent_bar(slide, support_x, y, support_w, Inches(0.06), C_BLUE)

        # Icon circle
        icons = ["Q", "S", "G"]  # Question, Sync, Group
        add_circle_badge(slide, support_x + Inches(0.2), y + Inches(0.25),
                        Inches(0.4), C_BLUE, icons[i], 12, C_WHITE)

        # Title
        add_ml_text(slide, support_x + Inches(0.75), y + Inches(0.2),
                    support_w - Inches(1.0), Inches(0.3),
                    title, size=15, color=C_BLUE, bold=True)
        # Subtitle
        add_ml_text(slide, support_x + Inches(0.75), y + Inches(0.5),
                    support_w - Inches(1.0), Inches(0.25),
                    subtitle, size=12, color=C_TEXT_DARK)
        # Description
        add_ml_text(slide, support_x + Inches(0.75), y + Inches(0.82),
                    support_w - Inches(1.0), Inches(0.7),
                    desc, size=10, color=C_TEXT_MUTED, line_spacing=1.35)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15: Getting Started -- 3-column visual guide with flow diagrams
# ═══════════════════════════════════════════════════════════════

def build_slide_15(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "上手准备",
                    "项目开始前，你需要了解的工具与流程 -- 一个下午就能搞定")

    col_w = Inches(3.55)
    total_w = 3 * col_w + 2 * Inches(0.35)
    start_x = (SLIDE_W - total_w) / 2
    cols_x = [start_x + i * (col_w + Inches(0.35)) for i in range(3)]
    col_start_y = Inches(1.75)
    col_h = Inches(5.3)

    # ═══ COLUMN 1: GitHub Workflow -- vertical flowchart ═══
    x1 = cols_x[0]
    add_accent_bar(slide, x1, col_start_y, col_w, Inches(0.42), C_BLUE)
    add_ml_text(slide, x1 + Inches(0.1), col_start_y + Inches(0.03),
                col_w - Inches(0.2), Inches(0.36),
                "GitHub 工作流", size=14, color=C_WHITE, bold=True,
                align=PP_ALIGN.CENTER)

    git_flow = [
        ("1", "Fork 仓库", "复制项目到你的账号下"),
        ("2", "Clone 到本地", "git clone 拉取代码"),
        ("3", "创建分支", "在开发分支上工作"),
        ("4", "编写代码", "实现你的编译器功能"),
        ("5", "提交 PR", "发起 Pull Request"),
        ("6", "Review 合入", "Code Review 后合并"),
    ]

    step_card_w = col_w - Inches(0.25)
    step_card_h = Inches(0.48)
    gf_start_y = col_start_y + Inches(0.55)

    for j, (num, label, desc) in enumerate(git_flow):
        y = gf_start_y + j * Inches(0.62)
        sx = x1 + Inches(0.125)

        add_rounded_card(slide, sx, y, step_card_w, step_card_h,
                        C_BLUE_LIGHT, C_BLUE, Pt(0.8))
        add_circle_badge(slide, sx + Inches(0.08), y + Inches(0.07),
                        Inches(0.3), C_BLUE, num, 9, C_WHITE)
        add_ml_text(slide, sx + Inches(0.48), y + Inches(0.03),
                    step_card_w - Inches(0.6), Inches(0.2),
                    label, size=12, color=C_BLUE, bold=True)
        add_ml_text(slide, sx + Inches(0.48), y + Inches(0.25),
                    step_card_w - Inches(0.6), Inches(0.2),
                    desc, size=8, color=C_TEXT_MUTED)

        # Down arrow
        if j < len(git_flow) - 1:
            add_arrow_down(slide, sx + step_card_w / 2 - Inches(0.05),
                          y + step_card_h, Inches(0.1), Inches(0.08), C_BLUE)

    # GitHub link note at bottom of column
    y_gh_note = gf_start_y + 6 * Inches(0.62) + Inches(0.05)
    add_ml_text(slide, x1 + Inches(0.15), y_gh_note,
                col_w - Inches(0.3), Inches(0.3),
                "查阅「上手文档」了解详细步骤",
                size=10, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ═══ COLUMN 2: AI-Assisted Development ═══
    x2 = cols_x[1]
    add_accent_bar(slide, x2, col_start_y, col_w, Inches(0.42), C_GOLD)
    add_ml_text(slide, x2 + Inches(0.1), col_start_y + Inches(0.03),
                col_w - Inches(0.2), Inches(0.36),
                "AI 辅助开发", size=14, color=C_WHITE, bold=True,
                align=PP_ALIGN.CENTER)

    ai_sections = [
        ("AI 编程助手", "豆包 & DeepSeek",
         "在线 AI 编程助手帮你快速\n验证想法，加速原型开发。\n不懂的代码？问 AI 解释。"),
        ("Agent 模式", "Claude Agent 体验",
         "对话即编程的全新体验。\n描述你想要的功能，\nAI 帮你生成、测试、迭代。"),
        ("加速学习", "AI 驱动的代码理解",
         "用 AI 解释陌生代码逻辑，\n阅读源码不再是障碍。\n学习效率翻倍，门槛降到最低。"),
    ]

    ai_card_h = Inches(1.35)
    ai_start_y = col_start_y + Inches(0.55)

    for j, (subtitle, name, desc) in enumerate(ai_sections):
        y = ai_start_y + j * (ai_card_h + Inches(0.1))
        sx = x2 + Inches(0.125)

        add_rounded_card(slide, sx, y, col_w - Inches(0.25), ai_card_h,
                        C_GOLD_LIGHT, C_GOLD, Pt(1))
        # Small top accent in each card
        add_accent_bar(slide, sx, y, col_w - Inches(0.25), Inches(0.04), C_GOLD)

        # Card number
        add_circle_badge(slide, sx + Inches(0.12), y + Inches(0.15),
                        Inches(0.3), C_GOLD, str(j + 1), 9, C_WHITE)
        # Subtitle
        add_ml_text(slide, sx + Inches(0.52), y + Inches(0.12),
                    col_w - Inches(0.8), Inches(0.22),
                    subtitle, size=12, color=C_GOLD, bold=True)
        # Name
        add_ml_text(slide, sx + Inches(0.52), y + Inches(0.36),
                    col_w - Inches(0.8), Inches(0.22),
                    name, size=10, color=C_TEXT_DARK, bold=True)
        # Description
        add_ml_text(slide, sx + Inches(0.52), y + Inches(0.6),
                    col_w - Inches(0.8), Inches(0.7),
                    desc, size=9, color=C_TEXT_MUTED, line_spacing=1.35)

    # Bottom note for AI column
    y_ai_note = ai_start_y + 3 * (ai_card_h + Inches(0.1))
    add_ml_text(slide, x2 + Inches(0.15), y_ai_note + Inches(0.05),
                col_w - Inches(0.3), Inches(0.3),
                "AI 是你最好的结对编程伙伴",
                size=10, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER, italic=True)

    # ═══ COLUMN 3: CI & Benchmark -- horizontal flow ═══
    x3 = cols_x[2]
    add_accent_bar(slide, x3, col_start_y, col_w, Inches(0.42), C_GREEN)
    add_ml_text(slide, x3 + Inches(0.1), col_start_y + Inches(0.03),
                col_w - Inches(0.2), Inches(0.36),
                "CI & Benchmark", size=14, color=C_WHITE, bold=True,
                align=PP_ALIGN.CENTER)

    ci_start_y = col_start_y + Inches(0.55)

    # CI flow -- each step is a compact card with arrow
    ci_steps = [
        ("提交代码", "git push 触发\nCI 流水线"),
        ("自动测试", "单元测试 +\n集成测试"),
        ("查看结果", "GitHub Actions\n实时反馈"),
        ("添加 Case", "为你的功能\n添加 CI 用例"),
        ("Benchmark", "LLVM vs ScratchV\n性能对比"),
    ]

    ci_card_h = Inches(0.62)
    for j, (label, desc) in enumerate(ci_steps):
        y = ci_start_y + j * (ci_card_h + Inches(0.08))
        sx = x3 + Inches(0.125)

        add_rounded_card(slide, sx, y, col_w - Inches(0.25), ci_card_h,
                        C_GREEN_LIGHT, C_GREEN, Pt(0.8))
        add_circle_badge(slide, sx + Inches(0.08), y + Inches(0.12),
                        Inches(0.32), C_GREEN, str(j + 1), 9, C_WHITE)
        add_ml_text(slide, sx + Inches(0.5), y + Inches(0.04),
                    col_w - Inches(0.65), Inches(0.22),
                    label, size=12, color=C_GREEN, bold=True)
        add_ml_text(slide, sx + Inches(0.5), y + Inches(0.28),
                    col_w - Inches(0.65), Inches(0.3),
                    desc, size=8, color=C_TEXT_MUTED, line_spacing=1.2)

        # Arrow
        if j < len(ci_steps) - 1:
            add_arrow_down(slide,
                          sx + (col_w - Inches(0.25)) / 2 - Inches(0.05),
                          y + ci_card_h, Inches(0.1), Inches(0.05), C_GREEN)

    # Bottom note for CI column
    y_ci_note = ci_start_y + 5 * (ci_card_h + Inches(0.08))
    add_ml_text(slide, x3 + Inches(0.15), y_ci_note + Inches(0.05),
                col_w - Inches(0.3), Inches(0.3),
                "每次提交自动运行，质量有保障",
                size=10, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER, italic=True)

    # ═══ Bottom: Learning path summary banner ═══
    banner_y = Inches(7.05)
    # Thin line across bottom
    add_section_divider(slide, MARGIN, banner_y - Inches(0.05), CONTENT_W, C_TEAL)
    add_ml_text(slide, MARGIN, banner_y, CONTENT_W, Inches(0.3),
                "学习路径：Fork 代码 -> AI 辅助开发 -> 提交 PR -> CI 自动验证 -> Benchmark 对比 -> 成就感满满！",
                size=12, color=C_TEAL, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 16: FAQ Part 1
# ═══════════════════════════════════════════════════════════════

def build_slide_16(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "常见疑问")

    qa_pairs = [
        ("Q: 没听过 ONNX，能行吗？",
         "A: 当然可以。第 1 周就会带你跑通例子，ONNX 只是一个文件格式，把它当成「模型存盘」就好。"),
        ("Q: 没学过编译原理，会不会听不懂？",
         "A: 我们会用直观的比喻（编译器就像「翻译官」），避开理论轰炸，先动手再做总结。"),
        ("Q: 需要买 RISC-V 开发板吗？",
         "A: 不需要。全程用软件模拟器（tinyfive），在笔记本电脑上就能跑。"),
        ("Q: 中途跟不上怎么办？",
         "A: 每个阶段有进度检查，mentor 会主动帮忙。有缓冲时间，可选只完成核心路径。完成比完美重要。"),
    ]

    q_bar_h = Inches(0.42)
    a_h = Inches(0.75)
    start_y = Inches(1.8)
    q_w = Inches(11.5)
    q_x = Inches(0.9)
    spacing = Inches(0.22)

    for i, (q, a) in enumerate(qa_pairs):
        y = start_y + i * (q_bar_h + a_h + spacing)
        # Q bar
        add_rounded_card(slide, q_x, y, q_w, q_bar_h, C_TEAL, None)
        add_ml_text(slide, q_x + Inches(0.2), y + Inches(0.03),
                    q_w - Inches(0.4), q_bar_h - Inches(0.04),
                    q, size=15, color=C_WHITE, bold=True)
        # A text
        add_ml_text(slide, q_x + Inches(0.35), y + q_bar_h + Inches(0.12),
                    q_w - Inches(0.7), a_h,
                    a, size=13, color=C_TEXT_DARK, line_spacing=1.35)


# ═══════════════════════════════════════════════════════════════
# SLIDE 17: FAQ Part 2
# ═══════════════════════════════════════════════════════════════

def build_slide_17(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    content_slide_decor(slide)
    add_slide_title(slide, "常见疑问（续）")

    qa_pairs = [
        ("Q: 结项后还想继续参与项目开发可以吗？",
         "A: ScratchV 开源社区会持续活跃，欢迎大家踊跃提交 PR！"),
        ("Q: 精选课题对我来说太难，或者我有其他自己想做的课题可以选吗？",
         "A: 精选课题只是提供一个 list 供大家选择，大家可以自由提出新的想法，"
         "与 mentor 讨论后立项并进行开发。"),
        ("Q: 我想做的课题有其他同学选择了，我还可以做吗？",
         "A: 我们鼓励合作完成课题，可以以小组为单位提交 PR，"
         "但是不能两个不同小组实现同一个课题。"),
        ("Q: 我想多锻炼自己！能不能多认领几个课题做？",
         "A: 每人在每个时刻只能认领最多一个课题，如果完成周期较短，"
         "可以在课题完成后立刻开始下一个课题。"),
    ]

    q_bar_h = Inches(0.42)
    start_y = Inches(1.8)
    q_w = Inches(11.5)
    q_x = Inches(0.9)

    for i, (q, a) in enumerate(qa_pairs):
        a_h = Inches(0.85) if len(a) > 60 else Inches(0.65)
        spacing = Inches(0.22)
        y = start_y + i * (q_bar_h + a_h + spacing)
        # Q bar
        add_rounded_card(slide, q_x, y, q_w, q_bar_h, C_TEAL, None)
        add_ml_text(slide, q_x + Inches(0.2), y + Inches(0.03),
                    q_w - Inches(0.4), q_bar_h - Inches(0.04),
                    q, size=15, color=C_WHITE, bold=True)
        # A text
        add_ml_text(slide, q_x + Inches(0.35), y + q_bar_h + Inches(0.12),
                    q_w - Inches(0.7), a_h,
                    a, size=13, color=C_TEXT_DARK, line_spacing=1.35)


# ═══════════════════════════════════════════════════════════════
# SLIDE 18: Call to Action -- with clickable link
# ═══════════════════════════════════════════════════════════════

def build_slide_18(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK_BG)

    # Decorative elements
    add_decor_ovals(slide, [(10.5, 5.0), (11.5, 4.0)], Inches(1.8))
    add_decor_ovals(slide, [(-1.0, 1.0)], Inches(1.8))
    # Accent dots
    for (dx, dy) in [(12.0, 0.5), (12.5, 6.5), (0.5, 6.0)]:
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(dx), Inches(dy), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_TEAL
        dot.line.fill.background()

    # Top accent
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), C_TEAL)

    # Main CTA
    add_ml_text(slide, Inches(1.2), Inches(0.8), Inches(10.0), Inches(1.0),
                "从今天起，给自己一个\n「创造编译器」的机会",
                size=34, color=C_TEXT_WHITE, bold=True, line_spacing=1.35)

    add_ml_text(slide, Inches(1.2), Inches(2.3), Inches(10.0), Inches(0.8),
                "也许你现在觉得编译器遥不可及，\n"
                "但三个月后，你会看着自己写的代码，把一行行模型规则变成芯片指令。",
                size=16, color=C_TEXT_LIGHT, line_spacing=1.45)

    add_ml_text(slide, Inches(1.2), Inches(3.4), Inches(10.0), Inches(0.5),
                "「我居然做到了」-- 这会是大学期间最难忘的回忆之一。",
                size=16, color=C_GOLD, italic=True)

    add_ml_text(slide, Inches(1.2), Inches(4.1), Inches(10.0), Inches(0.5),
                "不要让「基础不够」成为不敢开始的理由。",
                size=16, color=C_TEXT_WHITE, bold=True)

    # Clickable registration link
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(5.0), Inches(10.0), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    run_label = p.add_run()
    run_label.text = "立即报名: "
    run_label.font.size = Pt(18)
    run_label.font.color.rgb = C_TEXT_WHITE
    run_label.font.bold = True
    run_label.font.name = FONT

    run_url = p.add_run()
    run_url.text = "https://your-form-link.com"
    run_url.font.size = Pt(18)
    run_url.font.color.rgb = C_TEAL
    run_url.font.underline = True
    run_url.font.bold = True
    run_url.font.name = FONT
    run_url.hyperlink.address = "https://your-form-link.com"

    # QQ group
    add_ml_text(slide, Inches(1.2), Inches(5.6), Inches(10.0), Inches(0.4),
                "咨询（QQ 群）：1106852304",
                size=14, color=C_TEXT_LIGHT)

    # Footer
    add_accent_bar(slide, Inches(1.2), Inches(6.25), Inches(10.0), Inches(0.02), C_GOLD)
    add_ml_text(slide, Inches(1.2), Inches(6.35), Inches(10.0), Inches(0.5),
                "欢迎转发给同样好奇的小伙伴！\n"
                "你不需要很厉害才能开始，但你需要开始才能很厉害。",
                size=14, color=C_GOLD, align=PP_ALIGN.LEFT, line_spacing=1.4)


# ═══════════════════════════════════════════════════════════════
# SLIDE 19: Thank You / Closing
# ═══════════════════════════════════════════════════════════════

def build_slide_19(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK_BG)

    # Decorative
    add_decor_ovals(slide, [(10.5, 5.0), (11.5, 4.0)], Inches(1.8))
    add_decor_ovals(slide, [(-1.0, 1.0)], Inches(1.8))
    for (dx, dy) in [(12.0, 0.3), (0.5, 6.5)]:
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(dx), Inches(dy), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_TEAL
        dot.line.fill.background()

    # Top accent
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), C_TEAL)

    # Main thank you
    add_ml_text(slide, Inches(1.2), Inches(1.5), Inches(11.0), Inches(1.2),
                "Thank You", size=56, color=C_TEXT_WHITE, bold=True,
                align=PP_ALIGN.CENTER)

    # Gold accent under title
    add_accent_bar(slide, Inches(5.5), Inches(2.7), Inches(2.3), Inches(0.03), C_GOLD)

    add_ml_text(slide, Inches(1.2), Inches(3.2), Inches(11.0), Inches(0.8),
                "期待与你一起，开启编译器的奇妙旅程",
                size=22, color=C_TEAL, align=PP_ALIGN.CENTER)

    add_ml_text(slide, Inches(1.2), Inches(4.3), Inches(11.0), Inches(0.5),
                "Questions & Answers", size=18, color=C_TEXT_LIGHT,
                align=PP_ALIGN.CENTER)

    add_accent_bar(slide, Inches(5.0), Inches(5.1), Inches(3.3), Inches(0.02), C_TEXT_MUTED)
    add_ml_text(slide, Inches(1.2), Inches(5.4), Inches(11.0), Inches(0.5),
                "github.com/scratchv/scratchv   |   #ScratchV",
                size=12, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_ml_text(slide, Inches(1.2), Inches(6.2), Inches(11.0), Inches(0.5),
                "QQ 群：1106852304",
                size=13, color=C_TEXT_LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.500)

    # Remove default slide
    if prs.slides:
        xml_slides = prs.slides._sldIdLst
        for sldId in list(xml_slides):
            xml_slides.remove(sldId)

    # Build all 19 slides
    build_slide_01(prs)
    build_slide_02(prs)
    build_slide_03(prs)
    build_slide_04(prs)
    build_slide_05(prs)
    build_slide_06(prs)
    build_slide_07(prs)
    build_slide_08(prs)
    build_slide_09(prs)
    build_slide_10(prs)
    build_slide_11(prs)
    build_slide_12(prs)
    build_slide_13(prs)
    build_slide_14(prs)
    build_slide_15(prs)
    build_slide_16(prs)
    build_slide_17(prs)
    build_slide_18(prs)
    build_slide_19(prs)

    output = "/home/kinsomwang/workspace/ScratchV/ScratchV_宣讲PPT_polished.pptx"
    prs.save(output)
    print(f"Saved: {output}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
